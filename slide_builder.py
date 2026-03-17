"""Build Quantum Ethereal PPTX slides from script.json using python-pptx."""

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from validator import SETTINGS

# Design constants from config
COLORS = SETTINGS["design"]["colors"]
FONTS = SETTINGS["design"]["fonts"]
SLIDE_TYPES = SETTINGS["design"]["slide_types"]

# Dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)  # Standard 16:9
SLIDE_HEIGHT = Inches(7.5)

# Color palette
BG_COLOR = RGBColor.from_string(COLORS["bg"])
GOLD = RGBColor.from_string(COLORS["gold"])
GOLD_LIGHT = RGBColor.from_string(COLORS["gold_light"])
GOLD_DIM = RGBColor.from_string(COLORS["gold_dim"])
WHITE = RGBColor.from_string(COLORS["white"])
DIM = RGBColor.from_string(COLORS["dim"])
BLUE = RGBColor.from_string(COLORS["blue"])
VIOLET = RGBColor.from_string(COLORS["violet"])
TEAL = RGBColor.from_string(COLORS["teal"])
ROSE = RGBColor.from_string(COLORS["rose"])

BODY_WHITE = RGBColor(0xF0, 0xED, 0xE4)  # Warm white for max legibility
CITATION_GRAY = RGBColor(0xA8, 0xA4, 0xA0)  # Visible but secondary

ACCENT_COLORS = {
    "gold": GOLD,
    "blue": BLUE,
    "violet": VIOLET,
    "teal": TEAL,
    "rose": ROSE,
}

# Font names (with fallbacks)
FONT_DISPLAY = FONTS["fallback_display"]  # Georgia (safe fallback)
FONT_BODY = FONTS["fallback_body"]        # Calibri (safe fallback)

# Asset paths
ASSETS_DIR = Path(__file__).parent / "assets"
BG_MAP = {
    "title": "bg_title.png",
    "hook": "bg_title.png",
    "content": "bg_content.png",
    "science": "bg_science.png",
    "practice": "bg_practice.png",
    "quote": "bg_quote.png",
    "summary": "bg_summary.png",
    "cta": "bg_summary.png",
}
MOON_PATH = ASSETS_DIR / "decorations" / "moon_face.png"
DIVIDER_PATH = ASSETS_DIR / "decorations" / "divider_star.png"
CORNER_PATH = ASSETS_DIR / "decorations" / "corner_ornaments.png"
CORNER_PATHS = {
    "tl": ASSETS_DIR / "decorations" / "corner_tl.png",
    "tr": ASSETS_DIR / "decorations" / "corner_tr.png",
    "bl": ASSETS_DIR / "decorations" / "corner_bl.png",
    "br": ASSETS_DIR / "decorations" / "corner_br.png",
}
CONSTELLATION_PATH = ASSETS_DIR / "decorations" / "constellation_overlay.png"

# Layout constants — 10% padding on all edges
PAD_H = Inches(1.333)   # 10% of 13.333"
PAD_V = Inches(0.75)    # 10% of 7.5"
CONTENT_W = Inches(10.667)  # 13.333 - 2*1.333


def _set_slide_bg(slide, prs, slide_type="content"):
    """Set slide background: image if available, else solid color."""
    bg_file = BG_MAP.get(slide_type, "bg_content.png")
    bg_path = ASSETS_DIR / "backgrounds" / bg_file
    if bg_path.exists():
        slide.shapes.add_picture(
            str(bg_path), 0, 0, prs.slide_width, prs.slide_height
        )
    else:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR


def _add_text_box(slide, left, top, width, height, text, font_name=FONT_BODY,
                  font_size=Pt(18), font_color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with styled text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Set vertical anchor
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_separator_line(slide, left, top, width, color=GOLD_DIM):
    """Add a thin horizontal gold separator line."""
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, Pt(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_watermark(slide):
    """Add 'SELENE ACADEMIA' watermark at bottom — 14pt gold dim, wide spacing."""
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = "S E L E N E   A C A D E M I A"
    run.font.name = FONT_BODY
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x8B, 0x76, 0x35)  # Gold dim
    # Letter spacing via XML (300 = 3pt spacing)
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:spc'), '300')


def _set_picture_opacity(pic, opacity_pct):
    """Set opacity on a picture shape via XML (0-100)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    blipFill = pic._element.find(qn('p:blipFill'))
    if blipFill is not None:
        blip = blipFill.find(qn('a:blip'))
        if blip is not None:
            amt = int(opacity_pct * 1000)  # 40% = 40000
            alphaModFix = etree.SubElement(blip, qn('a:alphaModFix'))
            alphaModFix.set('amt', str(amt))


def _add_corner_ornaments(slide, color=GOLD_DIM):
    """Add top-2 corner ornaments only, ~18% opacity (barely-visible watermark)."""
    cw = Inches(1.4)
    ch = Inches(0.93)
    margin = Inches(0.15)
    sw = Inches(13.333)

    # Top corners only — bottom corners removed to avoid competing with content
    positions = {
        "tl": (margin, margin),
        "tr": (sw - cw - margin, margin),
    }

    has_images = all(CORNER_PATHS[k].exists() for k in positions)
    if has_images:
        for key, (left, top) in positions.items():
            pic = slide.shapes.add_picture(
                str(CORNER_PATHS[key]), left, top, cw, ch
            )
            _set_picture_opacity(pic, 18)
    else:
        for left, top in positions.values():
            _add_text_box(
                slide, left, top, Inches(0.5), Inches(0.4),
                text="✦",
                font_size=Pt(10), font_color=color,
                alignment=PP_ALIGN.CENTER
            )


def _add_moon(slide):
    """Add moon_face decoration on title slides — small, centered, subtle."""
    if MOON_PATH.exists():
        # ~150px wide ≈ 1.56". Moon is portrait (2:3), so height ≈ 2.34"
        moon_w = Inches(1.56)
        moon_h = Inches(2.34)
        # Center horizontally, with breathing room from top
        left = (Inches(13.333) - moon_w) / 2
        top = Inches(0.8)
        pic = slide.shapes.add_picture(
            str(MOON_PATH), int(left), int(top), moon_w, moon_h
        )
        _set_picture_opacity(pic, 60)  # Subtle, not protagonist


def _add_divider(slide, top=None):
    """Add star divider — 40% slide width, centered."""
    if DIVIDER_PATH.exists():
        div_w = Inches(5.333)  # 40% of 13.333"
        div_h = Inches(0.35)
        left = (Inches(13.333) - div_w) / 2
        if top is None:
            top = Inches(1.8)
        slide.shapes.add_picture(
            str(DIVIDER_PATH), int(left), int(top), div_w, div_h
        )


def _add_constellation(slide):
    """Add constellation overlay at low opacity for content/science slides.

    The PNG already has a transparent background (black removed).
    """
    if CONSTELLATION_PATH.exists():
        slide.shapes.add_picture(
            str(CONSTELLATION_PATH), 0, 0,
            Inches(13.333), Inches(7.5)
        )


def _add_bullet(slide, left, top, width, text, accent_color=GOLD):
    """Add a bullet with gold diamond and white body text for max contrast."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.75))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    # Set line spacing to 1.8
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': '180000'})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)

    # Gold diamond marker — proportional to bullet text
    dot_run = p.add_run()
    dot_run.text = "◆  "
    dot_run.font.name = FONT_BODY
    dot_run.font.size = Pt(24)
    dot_run.font.color.rgb = accent_color
    dot_run.font.bold = False
    # White body text — min 30pt per mobile legibility rule
    text_run = p.add_run()
    text_run.text = text
    text_run.font.name = FONT_BODY
    text_run.font.size = Pt(32)
    text_run.font.color.rgb = BODY_WHITE
    text_run.font.bold = False
    return txBox


def _add_citation(slide, left, top, width, citation_text):
    """Add a structured citation block — 16pt, italic, gray, split by | into lines."""
    from pptx.oxml.ns import qn

    # Split multiple citations by | separator
    citations = [c.strip() for c in citation_text.split("|")]

    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, cite in enumerate(citations):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER

        # Line spacing 1.5 for multi-line citations
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': '150000'})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)

        # Emoji on first line only
        prefix = "📚  " if idx == 0 else "      "
        icon_run = p.add_run()
        icon_run.text = prefix
        icon_run.font.name = FONT_BODY
        icon_run.font.size = Pt(18)
        icon_run.font.color.rgb = CITATION_GRAY

        # Citation text — italic, min 18pt for mobile
        text_run = p.add_run()
        text_run.text = cite
        text_run.font.name = FONT_BODY
        text_run.font.size = Pt(18)
        text_run.font.color.rgb = CITATION_GRAY
        text_run.font.italic = True

    return txBox


def _get_accent_color(slide_type: str) -> RGBColor:
    """Get the accent color for a slide type."""
    type_config = SLIDE_TYPES.get(slide_type, SLIDE_TYPES.get("content", {}))
    accent_name = type_config.get("accent", "gold")
    return ACCENT_COLORS.get(accent_name, GOLD)


def _get_icon(slide_type: str) -> str:
    """Get the icon for a slide type."""
    type_config = SLIDE_TYPES.get(slide_type, SLIDE_TYPES.get("content", {}))
    return type_config.get("icon", "◆")


def build_slide(prs, slide_data: dict, slide_number: int, total_slides: int):
    """Build a single slide from slide data."""
    slide_type = slide_data.get("type", "content")
    title = slide_data.get("title", "")
    bullets = slide_data.get("bullets", [])
    narration = slide_data.get("narration", "")
    citation = slide_data.get("citation")

    accent = _get_accent_color(slide_type)
    icon = _get_icon(slide_type)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, prs, slide_type)
    _add_corner_ornaments(slide, GOLD_DIM)
    _add_watermark(slide)

    # Slide type badge (top-left area) — skip for title slides
    if slide_type != "title":
        badge_text = f"{icon}  {slide_type.upper()}"
        _add_text_box(
            slide,
            left=PAD_H, top=PAD_V,
            width=Inches(3), height=Inches(0.35),
            text=badge_text,
            font_name=FONT_BODY, font_size=Pt(10),
            font_color=accent, bold=True
        )

    # Slide counter (top-right)
    _add_text_box(
        slide,
        left=Inches(13.333) - PAD_H - Inches(1.5), top=PAD_V,
        width=Inches(1.5), height=Inches(0.35),
        text=f"{slide_number}/{total_slides}",
        font_name=FONT_BODY, font_size=Pt(10),
        font_color=DIM,
        alignment=PP_ALIGN.RIGHT
    )

    # ── TITLE SLIDE ──────────────────────────────────────────────
    if slide_type == "title":
        _add_moon(slide)
        # Title — 54pt gold, centered vertically
        _add_text_box(
            slide,
            left=PAD_H, top=Inches(3.2),
            width=CONTENT_W, height=Inches(1.5),
            text=title,
            font_name=FONT_DISPLAY, font_size=Pt(54),
            font_color=GOLD_LIGHT, bold=True,
            alignment=PP_ALIGN.CENTER
        )
        _add_divider(slide, top=Inches(4.8))
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            _add_text_box(
                slide,
                left=PAD_H, top=Inches(5.3),
                width=CONTENT_W, height=Inches(0.6),
                text=subtitle,
                font_name=FONT_BODY, font_size=Pt(24),
                font_color=BODY_WHITE,
                alignment=PP_ALIGN.CENTER
            )

    # ── QUOTE SLIDE ─────────────────────────────────────────────
    # Minimal: just quote + author, lots of breathing room
    elif slide_type == "quote":
        _add_text_box(
            slide,
            left=Inches(1.8), top=Inches(2.0),
            width=Inches(9.733), height=Inches(3.0),
            text=f"❝ {title}",
            font_name=FONT_DISPLAY, font_size=Pt(36),
            font_color=BODY_WHITE, bold=False,
            alignment=PP_ALIGN.CENTER
        )
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            _add_text_box(
                slide,
                left=Inches(1.8), top=Inches(5.2),
                width=Inches(9.733), height=Inches(0.6),
                text=f"— {subtitle}",
                font_name=FONT_BODY, font_size=Pt(24),
                font_color=GOLD,
                alignment=PP_ALIGN.CENTER
            )

    # ── HOOK SLIDE ──────────────────────────────────────────────
    # Question is protagonist: 44pt centered, citation below
    elif slide_type == "hook":
        _add_text_box(
            slide,
            left=PAD_H, top=Inches(2.2),
            width=CONTENT_W, height=Inches(2.5),
            text=title,
            font_name=FONT_DISPLAY, font_size=Pt(44),
            font_color=BODY_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER
        )
        if citation:
            _add_divider(slide, top=Inches(5.0))
            _add_citation(slide, PAD_H, Inches(5.5), CONTENT_W, citation)

    # ── CONTENT / SCIENCE / PRACTICE / SUMMARY ─────────────────
    else:
        # Vertically center the content block
        area_top = Inches(1.1)
        area_bottom = Inches(6.4)
        area_h = area_bottom - area_top

        title_h = Inches(1.1)         # 54pt title
        sep_gap = Inches(0.35)
        bullet_spacing = Inches(0.95) # 32pt bullets need more room
        n_bullets = len(bullets)
        citation_h = Inches(1.5) if citation else 0

        block_h = title_h + sep_gap + (n_bullets * bullet_spacing) + citation_h
        content_top = area_top + max(0, (area_h - block_h) / 2)

        # Title — 54pt white bold (min size rule)
        _add_text_box(
            slide,
            left=PAD_H, top=int(content_top),
            width=CONTENT_W, height=Inches(1.1),
            text=title,
            font_name=FONT_DISPLAY, font_size=Pt(54),
            font_color=BODY_WHITE, bold=True
        )
        sep_top = int(content_top + title_h)
        _add_separator_line(slide, PAD_H, sep_top, CONTENT_W, accent)

        # Bullets
        bullet_start = int(content_top + title_h + sep_gap)
        for i, bullet in enumerate(bullets):
            _add_bullet(slide, PAD_H, int(bullet_start + Inches(i * 0.95)),
                        CONTENT_W, bullet, accent)

        # Citation block
        if citation:
            cite_div_top = int(bullet_start + n_bullets * bullet_spacing + Inches(0.15))
            _add_divider(slide, top=cite_div_top)
            cite_top = int(cite_div_top + Inches(0.45))
            _add_citation(slide, PAD_H, cite_top, CONTENT_W, citation)

    # Speaker notes = narration text
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = narration

    return slide


def build_presentation(script: dict) -> Presentation:
    """Build a full PPTX presentation from a script."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides = script.get("slides", [])
    total = len(slides)

    for i, slide_data in enumerate(slides):
        build_slide(prs, slide_data, i + 1, total)

    return prs


def build_from_script_file(script_path: str | Path, output_path: str | Path = None) -> Path:
    """Load a script.json and generate a PPTX file."""
    script_path = Path(script_path)
    with open(script_path) as f:
        script = json.load(f)

    if output_path is None:
        output_path = script_path.parent / "slides.pptx"
    else:
        output_path = Path(output_path)

    prs = build_presentation(script)
    prs.save(str(output_path))
    print(f"  📊 PPTX saved: {output_path} ({len(script['slides'])} slides)")
    return output_path


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python slide_builder.py <script.json> [output.pptx]")
        sys.exit(1)
    out = sys.argv[2] if len(sys.argv) > 2 else None
    build_from_script_file(sys.argv[1], out)
