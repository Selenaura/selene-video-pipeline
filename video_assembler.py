"""Assemble final MP4 video from slides (PNG) + audio (MP3).

Rendering pipeline:
  1. pptxtoimages (LibreOffice→PDF→PNG) for high-fidelity slide rendering
  2. Fallback: Pillow text extraction from python-pptx shapes

Assembly pipeline:
  1. ffmpeg concat demuxer (fast, no re-encoding of images)
  2. Fallback: moviepy (ImageClip + AudioFileClip)
"""

import json
import shutil
import subprocess
from pathlib import Path

from validator import SETTINGS

VIDEO_CONFIG = SETTINGS["video"]
RESOLUTION = tuple(int(x) for x in VIDEO_CONFIG["resolution"].split("x"))
FPS = 24
SILENCE_GAP = VIDEO_CONFIG["inter_slide_silence"]
MIN_SLIDE_SECS = VIDEO_CONFIG["min_slide_seconds"]
WPM = VIDEO_CONFIG["words_per_minute"]

FFMPEG = shutil.which("ffmpeg")


def _get_slide_duration(slide_data: dict, slide_timings: dict, index: int) -> float:
    """Get duration for a slide from multiple sources."""
    duration = slide_data.get("duration_seconds")
    if not duration:
        timing = slide_timings.get(index, {})
        duration = timing.get("duration_seconds")
    if not duration:
        narration = slide_data.get("narration", "")
        word_count = len(narration.split())
        duration = max(word_count / (WPM / 60) + SILENCE_GAP, MIN_SLIDE_SECS)
    return float(duration)


def render_slides_to_png(pptx_path: str | Path, output_dir: str | Path) -> list[Path]:
    """Render PPTX slides to PNG images. Try pptxtoimages first, fallback to Pillow."""
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Try pptxtoimages (LibreOffice → PDF → PNG) for high-fidelity rendering
    # This produces pixel-perfect slides but requires LibreOffice to work correctly
    try:
        from pptxtoimages.tools import PPTXToImageConverter
        temp_dir = str(output_dir / "_pptx_temp")
        converter = PPTXToImageConverter(
            pptx_path=str(pptx_path.resolve()),
            output_dir=str(output_dir.resolve()),
            output_format="png",
            temp_dir=temp_dir,
        )
        images = converter.convert()
        # Rename to our naming convention
        png_files = []
        for i, img_path in enumerate(sorted(Path(p) if isinstance(p, str) else p for p in images)):
            img_path = Path(img_path)
            target = output_dir / f"slide_{i:03d}.png"
            if img_path != target and img_path.exists():
                img_path.rename(target)
            png_files.append(target)
        if png_files:
            print(f"  🖼 Rendered via pptxtoimages (LibreOffice): {len(png_files)} slides")
            return png_files
    except ImportError:
        pass  # pptxtoimages not installed
    except Exception as e:
        print(f"  ⚠ pptxtoimages/LibreOffice unavailable ({type(e).__name__}), using Pillow renderer")

    # Fallback: Pillow-based text extraction from python-pptx
    return _render_slides_pillow(pptx_path, output_dir)


def _render_slides_pillow(pptx_path: Path, output_dir: Path) -> list[Path]:
    """Fallback: render PPTX slides to PNG using Pillow text extraction."""
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    png_files = []
    width, height = RESOLUTION

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (width, height), color=(10, 10, 15))
        draw = ImageDraw.Draw(img)

        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        y_pos = 80
        for j, text in enumerate(texts):
            if not text or text == "✦":
                continue

            if j == 0 and len(text) < 20:
                font_size, color = 22, (201, 168, 76)
            elif j <= 1:
                font_size, color, y_pos = 48, (232, 213, 160), 160
            else:
                font_size, color = 32, (240, 237, 228)
                y_pos += 20

            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
            except (OSError, IOError):
                font = ImageFont.load_default()

            max_w = width - 200
            words = text.split()
            lines = []
            cur = []
            for word in words:
                test = " ".join(cur + [word])
                bbox = draw.textbbox((0, 0), test, font=font)
                if bbox[2] - bbox[0] > max_w and cur:
                    lines.append(" ".join(cur))
                    cur = [word]
                else:
                    cur.append(word)
            if cur:
                lines.append(" ".join(cur))

            for line in lines:
                draw.text((100, y_pos), line, fill=color, font=font)
                y_pos += font_size + 10
            y_pos += 15

        try:
            wm_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
        except (OSError, IOError):
            wm_font = ImageFont.load_default()
        draw.text((width // 2 - 80, height - 40), "SELENE ACADEMIA",
                  fill=(64, 64, 69), font=wm_font)

        png_path = output_dir / f"slide_{i:03d}.png"
        img.save(str(png_path), "PNG")
        png_files.append(png_path)

    return png_files


def _assemble_ffmpeg(lesson_dir: Path, png_files: list[Path], slides: list[dict],
                     slide_timings: dict, dry_run: bool) -> Path:
    """Assemble video using ffmpeg concat demuxer — much faster than moviepy."""
    audio_dir = lesson_dir / "audio"
    output_path = lesson_dir / "lesson.mp4"

    # Build per-slide segments, then concat
    segments_dir = lesson_dir / "segments"
    segments_dir.mkdir(parents=True, exist_ok=True)
    segment_list = []

    for i, slide_data in enumerate(slides):
        if i >= len(png_files):
            break

        duration = _get_slide_duration(slide_data, slide_timings, i)
        png_path = png_files[i]
        audio_path = audio_dir / f"slide_{i:03d}.mp3"
        segment_path = segments_dir / f"seg_{i:03d}.mp4"

        # Build ffmpeg command for this segment
        cmd = ["ffmpeg", "-y"]

        # Input: loop image for duration
        cmd += ["-loop", "1", "-t", str(duration), "-i", str(png_path)]

        # Input: audio if available and valid
        has_audio = (audio_path.exists() and audio_path.stat().st_size > 2048
                     and not dry_run)
        if has_audio:
            cmd += ["-i", str(audio_path)]

        # Video encoding
        cmd += [
            "-vf", f"scale={RESOLUTION[0]}:{RESOLUTION[1]}:force_original_aspect_ratio=decrease,pad={RESOLUTION[0]}:{RESOLUTION[1]}:(ow-iw)/2:(oh-ih)/2",
            "-c:v", VIDEO_CONFIG["codec"],
            "-preset", "ultrafast",
            "-pix_fmt", "yuv420p",
            "-r", str(FPS),
        ]

        if has_audio:
            cmd += ["-c:a", "aac", "-b:a", VIDEO_CONFIG["audio_bitrate"], "-shortest"]
        else:
            # Generate silent audio track for concat compatibility
            cmd += ["-f", "lavfi", "-t", str(duration), "-i", "anullsrc=r=44100:cl=stereo"]
            # Re-specify inputs: image=0, silent_audio=1
            cmd = ["ffmpeg", "-y",
                   "-loop", "1", "-t", str(duration), "-i", str(png_path),
                   "-f", "lavfi", "-t", str(duration), "-i", "anullsrc=r=44100:cl=stereo",
                   "-vf", f"scale={RESOLUTION[0]}:{RESOLUTION[1]}:force_original_aspect_ratio=decrease,pad={RESOLUTION[0]}:{RESOLUTION[1]}:(ow-iw)/2:(oh-ih)/2",
                   "-c:v", VIDEO_CONFIG["codec"],
                   "-preset", "ultrafast",
                   "-pix_fmt", "yuv420p",
                   "-r", str(FPS),
                   "-c:a", "aac", "-b:a", VIDEO_CONFIG["audio_bitrate"],
                   "-shortest"]

        cmd.append(str(segment_path))

        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            raise RuntimeError(f"ffmpeg segment {i} failed: {result.stderr[-500:]}")

        segment_list.append(segment_path)

    # Concat all segments
    concat_file = segments_dir / "concat.txt"
    with open(concat_file, "w") as f:
        for seg in segment_list:
            f.write(f"file '{seg.resolve()}'\n")

    # Mix ambient music if available
    music_path = lesson_dir / "ambient.mp3"
    if not music_path.exists():
        music_path = Path("assets") / "ambient.mp3"

    concat_cmd = [
        "ffmpeg", "-y",
        "-f", "concat", "-safe", "0", "-i", str(concat_file),
    ]

    if music_path.exists() and not dry_run:
        concat_cmd += ["-stream_loop", "-1", "-i", str(music_path)]
        concat_cmd += [
            "-filter_complex",
            "[0:a][1:a]amix=inputs=2:duration=first:dropout_transition=2,volume=1.0[aout]",
            "-map", "0:v", "-map", "[aout]",
        ]
        print(f"  🎵 Ambient music mixed")
    else:
        concat_cmd += ["-c", "copy"]

    concat_cmd.append(str(output_path))

    result = subprocess.run(concat_cmd, capture_output=True, text=True, timeout=600)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg concat failed: {result.stderr[-500:]}")

    # Calculate total duration
    total = sum(_get_slide_duration(s, slide_timings, i) for i, s in enumerate(slides) if i < len(png_files))
    print(f"  🎥 Video saved: {output_path} ({total:.0f}s / {total/60:.1f} min)")

    # Cleanup segments
    shutil.rmtree(segments_dir, ignore_errors=True)

    return output_path


def _assemble_moviepy(lesson_dir: Path, png_files: list[Path], slides: list[dict],
                      slide_timings: dict, dry_run: bool) -> Path:
    """Fallback: assemble video using moviepy."""
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips

    audio_dir = lesson_dir / "audio"
    clips = []

    for i, slide_data in enumerate(slides):
        if i >= len(png_files):
            break

        duration = _get_slide_duration(slide_data, slide_timings, i)
        png_path = png_files[i]
        audio_path = audio_dir / f"slide_{i:03d}.mp3"

        img_clip = ImageClip(str(png_path), duration=duration)
        img_clip = img_clip.resized(RESOLUTION)

        if (audio_path.exists() and audio_path.stat().st_size > 2048
                and not dry_run):
            try:
                audio_clip = AudioFileClip(str(audio_path))
                img_clip = img_clip.with_audio(audio_clip)
            except Exception as e:
                print(f"  ⚠ Could not load audio for slide {i}: {e}")

        clips.append(img_clip)

    if not clips:
        raise RuntimeError("No clips to assemble")

    print(f"  🎬 Assembling {len(clips)} clips (moviepy)...")
    final = concatenate_videoclips(clips, method="compose")

    # Ambient music
    music_path = lesson_dir / "ambient.mp3"
    if not music_path.exists():
        music_path = Path("assets") / "ambient.mp3"
    if music_path.exists() and not dry_run:
        try:
            music = AudioFileClip(str(music_path))
            if music.duration < final.duration:
                from moviepy import concatenate_audioclips
                loops_needed = int(final.duration / music.duration) + 1
                music = concatenate_audioclips([music] * loops_needed)
            music = music.subclipped(0, final.duration)
            music = music.with_volume_scaled(0.08)
            if final.audio is not None:
                from moviepy import CompositeAudioClip
                final = final.with_audio(CompositeAudioClip([final.audio, music]))
            else:
                final = final.with_audio(music)
            print(f"  🎵 Ambient music mixed at 8% volume")
        except Exception as e:
            print(f"  ⚠ Could not mix ambient music: {e}")

    output_path = lesson_dir / "lesson.mp4"
    final.write_videofile(
        str(output_path),
        fps=FPS,
        codec=VIDEO_CONFIG["codec"],
        preset=VIDEO_CONFIG["preset"],
        audio_codec="aac",
        audio_bitrate=VIDEO_CONFIG["audio_bitrate"],
        logger=None,
    )

    total_duration = sum(c.duration for c in clips)
    print(f"  🎥 Video saved: {output_path} ({total_duration:.0f}s / {total_duration/60:.1f} min)")

    final.close()
    for clip in clips:
        clip.close()

    return output_path


def assemble_video(lesson_dir: str | Path, dry_run: bool = False) -> Path:
    """Assemble final MP4 from slides + audio. Uses ffmpeg if available, else moviepy."""
    lesson_dir = Path(lesson_dir)
    script_path = lesson_dir / "script.json"
    pptx_path = lesson_dir / "slides.pptx"
    audio_dir = lesson_dir / "audio"
    manifest_path = audio_dir / "manifest.json"

    with open(script_path) as f:
        script = json.load(f)

    manifest = {}
    if manifest_path.exists():
        with open(manifest_path) as f:
            manifest = json.load(f)

    slide_timings = {s["slide_index"]: s for s in manifest.get("slides", [])}

    # Step 1: Export slides to PNG
    frames_dir = lesson_dir / "frames"
    print(f"  🖼 Exporting slides to PNG...")
    png_files = render_slides_to_png(pptx_path, frames_dir)
    print(f"  🖼 {len(png_files)} slide frames exported")

    slides = script.get("slides", [])

    # Step 2: Assemble — prefer ffmpeg for speed
    if FFMPEG:
        print(f"  🎬 Assembling {len(slides)} slides with ffmpeg...")
        try:
            return _assemble_ffmpeg(lesson_dir, png_files, slides, slide_timings, dry_run)
        except Exception as e:
            print(f"  ⚠ ffmpeg assembly failed ({e}), falling back to moviepy")

    return _assemble_moviepy(lesson_dir, png_files, slides, slide_timings, dry_run)


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python video_assembler.py <lesson_dir> [--dry-run]")
        sys.exit(1)
    dry = "--dry-run" in sys.argv
    assemble_video(sys.argv[1], dry_run=dry)
